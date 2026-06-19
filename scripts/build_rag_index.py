from __future__ import annotations
import json
import os
import sys
from pathlib import Path

import chromadb
from chromadb.utils.embedding_functions import OpenAIEmbeddingFunction
from pypdf import PdfReader
from pptx import Presentation

REPO_ROOT = Path(__file__).parent.parent
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))
from api.bootstrap_env import load_repo_dotenv

load_repo_dotenv()

DATA_DIR = REPO_ROOT / "data_extract"
# No need to specify CHROMA_PATH here — it's set by sync-and-ingest.sh script to a local temp dir on Cloud Run, and does not affect local development
CHROMA_PATH = Path(os.environ.get("CHROMA_PATH", str(DATA_DIR / "chroma_db")))
MANIFEST_PATH = DATA_DIR / "rag_manifest.json"

_openai_key = os.environ.get("OPENAI_API_KEY")
if not _openai_key:
    raise SystemExit(
        "OPENAI_API_KEY is not set. Add it to the repo .env (next to pyproject.toml) or export it, "
        "then re-run: uv run python scripts/build_rag_index.py"
    )

_EF = OpenAIEmbeddingFunction(
    api_key=_openai_key,
    model_name="text-embedding-3-small",
)


def extract_pdf(path: Path, category: str) -> list[dict]:
    reader = PdfReader(str(path))
    chunks = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if len(text.strip()) < 50:
            continue
        chunks.append({
            "text": text,
            "source": path.name,
            "category": category,
            "page": i,
        })
    return chunks


def extract_pptx(path: Path, category: str) -> list[dict]:
    prs = Presentation(str(path))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text if slide.shapes.title else ""
        body = " ".join(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape != slide.shapes.title
        )
        text = f"{title}\n{body}".strip()
        if len(text) < 30:
            continue
        chunks.append({
            "text": text,
            "source": path.name,
            "category": category,
            "page": i,
        })
    return chunks


def main() -> None:
    manifest = json.loads(MANIFEST_PATH.read_text())

    client = chromadb.PersistentClient(path=str(CHROMA_PATH))
    collection = client.get_or_create_collection("neom_docs", embedding_function=_EF)

    total_chunks = 0
    total_docs = 0

    for entry in manifest["documents"]:
        rel_path = entry["path"]
        category = entry["category"]
        full_path = DATA_DIR / rel_path

        if not full_path.exists():
            print(f"WARNING: file not found, skipping — {rel_path}")
            continue

        suffix = full_path.suffix.lower()
        if suffix == ".pdf":
            chunks = extract_pdf(full_path, category)
        elif suffix == ".pptx":
            chunks = extract_pptx(full_path, category)
        else:
            print(f"WARNING: unsupported file type, skipping — {rel_path}")
            continue

        if not chunks:
            print(f"WARNING: no usable chunks extracted — {rel_path}")
            continue

        ids = [f"{full_path.name}_{c['page']}" for c in chunks]
        texts = [c["text"] for c in chunks]
        metadatas = [{"source": c["source"], "category": c["category"], "page": c["page"]} for c in chunks]

        # upsert is idempotent — safe to re-run
        collection.upsert(ids=ids, documents=texts, metadatas=metadatas)

        total_chunks += len(chunks)
        total_docs += 1
        print(f"  {len(chunks):4d} chunks  ← {full_path.name}")

    print(f"\nIndexed {total_chunks} chunks from {total_docs} documents")


if __name__ == "__main__":
    main()
